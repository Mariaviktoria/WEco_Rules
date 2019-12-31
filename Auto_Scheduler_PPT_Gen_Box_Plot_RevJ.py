import pandas as pd
from matplotlib import pyplot as plt
from pptx.dml.color import RGBColor
#from pd2ppt import df_to_table
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from pandas.plotting import table
import numpy as np
from numpy import sqrt,mean,log,diff
import seaborn as sns
from pptx.util import Pt
import matplotlib.patches as mpatches
import os
import time
import schedule
from itertools import groupby, cycle

# from sys import exit


def data_list(csv_file):
    global dfcpk
    # Slicing the CPK, Parameters and the Spec limits cells
    dfcpk = pd.read_csv(csv_file, sep=',', index_col=False, header=None)
    dfcpk = dfcpk.iloc[:, 0:12]
    dfcpk.columns = ['Parameter', 'Info', 'Units', 'LSL', 'Target', 'USL', 'N', 'Mean', '3Sigma', 'Cp', 'K', 'Cpk']
    dfcpk1 = dfcpk.loc[:, 'LSL':].apply(pd.to_numeric, errors='coerce')
    dfcpk2 = dfcpk.loc[:, 'Parameter':'Units']
    dfcpk = pd.concat([dfcpk1, dfcpk2], axis=1)
    dfcpk.to_csv(logs + 'pcap_list.csv') # Generate list file with headers added from raw csv file
    global dftable
    dftable = dfcpk.iloc[:,:]
    cols = list(dftable.columns)
    cols = [cols[-1]] + cols[:-1]
    dftable = dftable[cols]
    dftable = dftable.round(2)
    def func(row):
        if row['Cpk'] >= 2.00:
            val = 'S'
        elif row['Cpk'] >= 1.67 and row['Cpk'] <2.00:
            val = 'E'
        elif row['Cpk'] >= 1.33 and row['Cpk'] <1.67:
            val = 'G'
        elif row['Cpk'] <1.33:
            val = 'N'
        else:
            val = 'NO GRADE AVAILABLE'
        return val
    dftable['Cpk Grade'] = dftable.apply(func, axis=1)

    #Count grade and calculate percentage
    grade_count_s = len(dftable.loc[dftable['Cpk Grade'] == 'S'])
    grade_count_e = len(dftable.loc[dftable['Cpk Grade'] == 'E'])
    grade_count_g = len(dftable.loc[dftable['Cpk Grade'] == 'G'])
    grade_count_n = len(dftable.loc[dftable['Cpk Grade'] == 'N'])
    grade_count_nga = len(dftable.loc[dftable['Cpk Grade'] == 'NO GRADE AVAILABLE'])
    global count_row, result, E, G, N, S
    # print(dftable.shape[0])
    # print(grade_count_nga)
    count_row = (dftable.shape[0] - grade_count_nga)
    # print(count_row)
    E = round(((grade_count_e/count_row)*100),2)
    G = round(((grade_count_g/count_row)*100),2)
    N = round(((grade_count_n/count_row)*100),2)
    S = round(((grade_count_s/count_row)*100),2)
    # print(grade_count[0])
    # print(grade_count[1])
    # print(grade_count[2])
    # print(grade_count[3])
    E = str('E='+str(E) + '%')
    G = str('G='+str(G) + '%')
    N = str('N='+str(N) + '%')
    S = str('S='+str(S) + '%')

def data_lot(txt_file):

    global file_name
    file_name = txt_file.split(".")[0]
    col_names = dfcpk.loc[:, 'Parameter']
    df1 = pd.read_csv(txt_file, delim_whitespace=True, names=col_names)
    df1.to_csv(logs + 'temp.csv', index=True)
    df2 = pd.read_csv(logs + 'temp.csv')
    df1 = df2.rename(columns={"Unnamed: 0": "lot", "Unnamed: 1": "product_name", "Unnamed: 2": "date", "Unnamed: 3": "wafers", "Unnamed: 4": "sites"})
    # print(df1.head())
    df3 = df1.loc[:, 'wafers':].apply(pd.to_numeric, errors='coerce')
    df4 = df1.loc[:, 'lot':'date']
    #for slicing lot numbers to display in drop down
    dflot = df1.loc[:, 'lot']
    dfprod = df1.loc[:, 'product_name']
    global df5, min_date, max_date, df6
    df5 = pd.concat([df3, df4], axis=1)
    df5['date'] = pd.to_datetime(df5['date'], format='%d-%m-%Y', errors='coerce')
    # print(df5['date'])
    min_date = (df5['date'].min()).date()
    # print(min_date)
    max_date = (df5['date'].max()).date()
    # print(max_date)
    # print(df5.head())
    df6 = df5.sort_values(by='date', ascending=True) #Sorting by date from latest to old
    df6.to_csv(logs + 'sort_by_dates.csv')
    df6['date'] = df6['date'].dt.date
    df6['date'] = df6['date'].astype(str)
    global temp
    temp = df6.groupby(['date', 'lot']).mean().reset_index()
    temp = temp.iloc[:, 0:2]

def cpk_select():

    global prs
    prs = Presentation()

    #Adding the results slide as the first slide
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'CPK REPORT'
    rows = 5
    cols = 2
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)
    # write column headings
    table.cell(0, 0).text = 'CPK GRADE'
    table.cell(0, 1).text = 'PERCENTAGE'
    # write body cells
    table.cell(2, 0).text = 'Excellent: 1.67<=Cpk<2.00'
    table.cell(2, 1).text = E
    table.cell(1, 0).text = 'Superior: 2.00<=Cpk'
    table.cell(1, 1).text = S
    table.cell(3, 0).text = 'Good: 1.33<=Cpk<1.67'
    table.cell(3, 1).text = G
    table.cell(4, 0).text = 'Action Required: Cpk<1.33'
    table.cell(4, 1).text = N

    # Change the below line to plot all parameters
    yval = dfcpk.iloc[0:3, 0:]
    Graphing(yval)
    prs.save(logs + file_name +'%d.pptx' %time.time())

def Graphing(yval):
    for i in yval.index:
        parm = yval.loc[i, 'Parameter']
        targetval = round((yval.loc[i, 'Target']), 2)
        minval = round((yval.loc[i, 'LSL']), 2)
        maxval = round((yval.loc[i, 'USL']), 2)
        cpkrange = round((yval.loc[i, 'Cpk']), 5)

        if cpkrange <1:
            border_color = RGBColor(0xFF, 0x00, 0x00)
        elif cpkrange >=1 and cpkrange <=1.3:
            border_color = RGBColor(0x00, 0x00, 0xFF)
        elif cpkrange >1.3:
            border_color = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            border_color = RGBColor(0x00, 0x00, 0x00)

        cp = round((yval.loc[i, 'Cp']), 5).astype(str)
        info = (yval.loc[i, 'Info'])
        units = (yval.loc[i, 'Units'])
        fig = plt.figure(figsize=(35, 15))
        ax = plt.subplot(111)
        # my_pal = {"0F2U8": "blue", "0F3H0": "orange", "0F3G9": "green", "0F3H1": "purple"}
        sns.boxplot(y=parm, x="lot", hue="product_name", data=df6, dodge=False, notch=True, palette=["blue", "orange", "green", "purple", "red", "pink", "yellow", "brown"])
        sns.pointplot(y=parm, x='lot', data=df6, linestyles='-', scale=0.6,
                      color='yellow', errwidth=0, capsize=0)
        plt.axhline(y=targetval, xmin=0, xmax=100, color='k', linewidth=2, linestyle='--', label='Target')
        plt.axhline(y=minval, xmin=0, xmax=100, color='r', linewidth=2, linestyle='--', label='LSL')
        plt.axhline(y=maxval, xmin=0, xmax=100, color='r', linewidth=2, linestyle='--', label='USL')
        plt.xticks(rotation=90, fontsize=12)
        plt.xlabel('Lot', fontsize=14)
        plt.ylabel(parm+units, fontsize=18)
        plt.title(str(min_date) + '~' + str(max_date) + "        " + str(info) + '  \nCp=' + str(cp) + '   Cpk=' +
                  str(            cpkrange) + '   Target=' + targetval.astype(str) + ' USL=' + maxval.astype(str) + ' LSL=' + minval.astype(str),
            fontsize=18)
        plt.suptitle("")
        # plt.legend()
        # fig.tight_layout()
        patch = mpatches.Patch(color='white', label=temp.values)
        ax.legend(handles = [patch],  bbox_to_anchor=(1, 1), loc=2, borderaxespad=0., prop={'size': 8})
        plt.tick_params(axis='y', which='major', labelsize=18,colors='black')
        fig.legend(loc=4, fontsize=9)
        plt.tight_layout()
        # plt.show()

        #Slide heading
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.font = Pt(15)
        title_shape.text = parm

        #modify image name to strip /
        parm = parm.replace("/", "_")
        plt.savefig(logs + str(cpkrange) + '_' + 'CPK' + '_' + parm + '.png', dpi=300)
        image = logs + str(cpkrange) + '_' + 'CPK' + '_' + parm + '.png'

        # title_slide_layout = prs.slide_layouts[6]
        # slide = prs.slides.add_slide(title_slide_layout)
        pic = slide.shapes.add_picture(image, Inches(.4), Inches(1.5), Inches(9.4), Inches(5.4))
        line = pic.line
        line.color.rgb = border_color
        line.width = Inches(0.05)
        plt.clf()
        plt.cla()
        plt.close('all')

def spc_control_rule():

    global gumc07, gumc08
    count_gwe2 = 0
    count_gwe3 = 0
    count_gwe4 = 0
    count_gwe6 = 0
    count_gwe7 = 0
    count_gwe8 = 0

    lot_num_gwe1 = []
    lot_num_gwe2 = []
    lot_num_gwe3 = []
    lot_num_gwe4 = []
    lot_num_gwe5 = []
    lot_num_gwe6 = []
    lot_num_gwe7 = []
    lot_num_gwe8 = []
    lot_num_gumc14 = []

    yval = dfcpk.iloc[0:, 0:]
    yval.to_csv(logs + 'yval.csv')
    lot_mean = df6.groupby(['lot']).mean().reset_index()
    lot_mean.to_csv(logs + 'lot_mean.csv')

    def CHECK_GUMC07(l, gumc07):
        list_gumc07 = lot_mean.iloc[:, [0, l + 3]]
        for k in range(len(list_gumc07) - 7):
            grp_7 = list_gumc07[k:7+k]
            p_value = grp_7.iloc[:, 1]
            lot = grp_7.iloc[:, 0]
            result = all(i < j for i, j in zip(p_value, p_value[1:]))
            if result is True:
                gumc07 = gumc07.append(lot).drop_duplicates().transpose()

    def CHECK_GUMC08(l, gumc08):
        list_gumc08 = lot_mean.iloc[:, [0, l + 3]]
        for k in range(len(list_gumc08) - 7):
            grp_7 = list_gumc08[k:7+k]
            p_value = grp_7.iloc[:, 1]
            lot = grp_7.iloc[:, 0]
            result = all(i > j for i, j in zip(p_value, p_value[1:]))
            if result is True:
                gumc08 = gumc08.append(lot).drop_duplicates().transpose()

    for i in yval.index:

        gwe1 = pd.DataFrame(columns=range(1))
        gwe2 = pd.DataFrame(columns=range(1))
        gwe3 = pd.DataFrame(columns=range(1))
        gwe4 = pd.DataFrame(columns=range(1))
        gwe5 = pd.DataFrame(columns=range(1))
        gwe6 = pd.DataFrame(columns=range(1))
        gwe7 = pd.DataFrame(columns=range(1))
        gwe8 = pd.DataFrame(columns=range(1))
        gumc14 = pd.DataFrame(columns=range(1))

        minval = round((yval.loc[i, 'LSL']), 2)
        maxval = round((yval.loc[i, 'USL']), 2)
        targetval = round((yval.loc[i, 'Target']), 2)
        three_sigma = round((yval.loc[i, '3Sigma']), 2)
        two_sigma = (three_sigma / 3) * 2
        one_sigma = three_sigma / 3
        parm = yval.loc[i, 'Parameter']
        parm = parm.replace("/", "_")

        # GUMC07 Rule Check
        gumc07 = pd.DataFrame(columns=range(1))
        CHECK_GUMC07(i, gumc07)

        # GUMC08 Rule Check
        gumc08 = pd.DataFrame(columns=range(1))
        CHECK_GUMC08(i, gumc08)


        #All Other Rule Check
        for j in range(0, len(lot_mean)):

            # GWE1 Rule Check
            if (round(lot_mean.iloc[j, i + 3], 2)) > (targetval + three_sigma):
                lot_num_gwe1.append(lot_mean.iloc[j, 0])
                gwe1 = gwe1.append(lot_num_gwe1, ignore_index=True).drop_duplicates()

            # GWE2 Rule Check
            if (round(lot_mean.iloc[j, i + 3], 2)) > (targetval + two_sigma):
                lot_num_gwe2.append(lot_mean.iloc[j, 0])
                count_gwe2 = count_gwe2 + 1
                if count_gwe2 == 2:
                    gwe2 = gwe2.append(lot_num_gwe2, ignore_index=True).drop_duplicates()
                    count_gwe2 = 0
                    lot_num_gwe2.clear()
            else:
                count_gwe2 = 0
                lot_num_gwe2.clear()

            # GWE3 Rule Check
            if (round(lot_mean.iloc[j, i + 3], 2)) > (targetval + one_sigma):
                lot_num_gwe3.append(lot_mean.iloc[j, 0])
                count_gwe3 = count_gwe3 + 1
                if count_gwe3 == 4:
                    gwe3 = gwe3.append(lot_num_gwe3, ignore_index=True).drop_duplicates()
                    count_gwe3 = 0
                    lot_num_gwe3.clear()
            else:
                count_gwe3 = 0
                lot_num_gwe3.clear()

            # GWE4 Rule Check
            if (round(lot_mean.iloc[j, i + 3], 2)) > targetval:
                lot_num_gwe4.append(lot_mean.iloc[j, 0])
                count_gwe4 = count_gwe4 + 1
                if count_gwe4 == 8:
                    gwe4 = gwe4.append(lot_num_gwe4, ignore_index=True).drop_duplicates()
                    count_gwe4 = 0
                    lot_num_gwe4.clear()
            else:
                count_gwe4 = 0
                lot_num_gwe4.clear()

            # GWE5 Rule Check
            if (round(lot_mean.iloc[j, i + 3], 2)) < (targetval - three_sigma):
                lot_num_gwe5.append(lot_mean.iloc[j, 0])
                gwe5 = gwe5.append(lot_num_gwe5, ignore_index=True).drop_duplicates()

            # GWE6 Rule Check
            if (round(lot_mean.iloc[j, i + 3], 2)) < (targetval - two_sigma):
                lot_num_gwe6.append(lot_mean.iloc[j, 0])
                count_gwe6 = count_gwe6 + 1
                if count_gwe6 == 2:
                    gwe6 = gwe6.append(lot_num_gwe6, ignore_index=True).drop_duplicates()
                    count_gwe6 = 0
                    lot_num_gwe6.clear()
            else:
                count_gwe6 = 0
                lot_num_gwe6.clear()

            # GWE7 Rule Check
            if (round(lot_mean.iloc[j, i + 3], 2)) < (targetval - one_sigma):
                lot_num_gwe7.append(lot_mean.iloc[j, 0])
                count_gwe7 = count_gwe7 + 1
                if count_gwe7 == 4:
                    gwe7 = gwe7.append(lot_num_gwe7, ignore_index=True).drop_duplicates()
                    count_gwe7 = 0
                    lot_num_gwe7.clear()
            else:
                count_gwe7 = 0
                lot_num_gwe7.clear()

            # GWE8 Rule Check
            if (round(lot_mean.iloc[j, i + 3], 2)) < targetval:
                lot_num_gwe8.append(lot_mean.iloc[j, 0])
                count_gwe8 = count_gwe8 + 1
                if count_gwe8 == 8:
                    gwe8 = gwe8.append(lot_num_gwe8, ignore_index=True).drop_duplicates()
                    count_gwe8 = 0
                    lot_num_gwe8.clear()
            else:
                count_gwe8 = 0
                lot_num_gwe8.clear()

            # GUMC14 Rule Check
            if not (round(lot_mean.iloc[j, i + 3], 2)) > minval and (round(lot_mean.iloc[j, i + 3], 2)) < maxval:
                lot_num_gumc14.append(lot_mean.iloc[j, 0])
                gumc14 = gumc14.append(lot_num_gumc14, ignore_index=True).drop_duplicates()


        gwe1 = gwe1.rename(columns={0: 'GWE1'})
        gwe2 = gwe2.rename(columns={0: 'GWE2'})
        gwe3 = gwe3.rename(columns={0: 'GWE3'})
        gwe4 = gwe4.rename(columns={0: 'GWE4'})
        gwe5 = gwe5.rename(columns={0: 'GWE5'})
        gwe6 = gwe6.rename(columns={0: 'GWE6'})
        gwe7 = gwe7.rename(columns={0: 'GWE7'})
        gwe8 = gwe8.rename(columns={0: 'GWE8'})
        gumc07 = gumc07.rename(columns={0: 'GUMC07'})
        gumc08 = gumc08.rename(columns={0: 'GUMC08'})
        gumc14 = gumc14.rename(columns={0: 'GUMC14'})

        df_mining = pd.concat([gwe1, gwe2, gwe3, gwe4, gwe5, gwe6, gwe7, gwe8, gumc07, gumc08, gumc14], axis=1)

        if not df_mining.empty:
            df_mining.to_csv(logs + parm + '.csv')
            # print(df_mining)


def trend_graph():
    global Cpk1, Cpk2, Cpk3, prs, top, left, right, width, height, path, logs

    path = 'C:/test/'
    logs = 'C:/test/logs/'

    folder = os.fsencode(path)

    raw_txt = []
    raw_csv = []

    #Iterate .csv files in directory
    for file in os.listdir(folder):
        filename = os.fsdecode(file)
        if filename.endswith('.csv'):  # whatever file types you're using...
            raw_csv.append(filename)

    #Iterate .txt files in directory
    for file in os.listdir(folder):
        filename = os.fsdecode(file)
        if filename.endswith('.txt'):  # whatever file types you're using...
            raw_txt.append(filename)

    list_size = len(raw_csv)
    #Loop thru all the iterated files to plot graphs
    for i in range(list_size):
        csv_file = raw_csv[i]
        data_list(csv_file)
        txt_file = raw_txt[i]
        data_lot(txt_file)
        spc_control_rule()
        cpk_select()

schedule.every(0).minutes.do(trend_graph) #This is for test purpose only. Uncommment the below line for actual scheduler
# schedule.every().tuesday.at("17:14").do(trend_graph)

while True:
    # Checks whether a scheduled task
    # is pending to run or not
    schedule.run_pending()
    time.sleep(0)



